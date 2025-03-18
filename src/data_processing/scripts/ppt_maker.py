from pptx import Presentation
from pptx.util import Inches

# Create a PowerPoint presentation
prs = Presentation()


run_nos = [426 + i for i in range(9)]
images = ['plot_XvsY.png', 'Drift1.png', 'Drift2.png', 'Drift3.png']
path_type = ["straight line", "straight line", "Closed Loop", "Closed Loop", "Figure 8", "Figure 8", "Closed Loop", "Closed Loop", "Arbitrary Path"]

# Loop through runs and plots
i = 0

for run in run_nos:
    folder_path =  f"/home/catkin_ws/src/data_processing/data/online_plotting/run_no{run}/"
    i+=1
    plot_num = 1
    for image in images:
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank slide layout
        title = slide.shapes.title
        print(len(path_type), i)
        title.text = f"Run: {i} " + path_type[i-1]
        # Add each plot as an image to the slide
        img_path = folder_path + image
        left = Inches(0.5 + (plot_num - 1) * 2)  # Adjust position dynamically
        top = Inches(1.5)  # Fixed vertical position
        slide.shapes.add_picture(img_path, left, top, width=Inches(10))

# Save the PowerPoint
prs.save("presentation.pptx")
